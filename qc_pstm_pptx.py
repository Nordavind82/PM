#!/usr/bin/env python3
"""
PSTM Migration QC Report Generator.

Creates PowerPoint presentations with QC visualizations for each offset bin:
- 3 inline sections
- 3 crossline sections
- 3 time slices (300, 500, 700 ms)

Also extracts and visualizes CIGs at selected locations.

Usage:
    python qc_pstm_pptx.py [--bins 0-37] [--output-dir PATH]
"""

import argparse
import sys
from datetime import datetime
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import numpy as np
import polars as pl
import zarr

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


# =============================================================================
# Configuration
# =============================================================================

MIGRATION_DIR = Path("/Volumes/AO_DISK/PSTM_common_offset")
INPUT_DIR = Path("/Users/olegadamovich/SeismicData/common_offset_gathers_new")
OUTPUT_DIR = Path("/Volumes/AO_DISK/PSTM_common_offset/qc_reports")

# Time slices to visualize
TIME_SLICES_MS = [300, 500, 700]

# Inline/crossline positions (evenly distributed)
INLINE_POSITIONS = [100, 255, 410]  # Near start, middle, near end
XLINE_POSITIONS = [85, 213, 340]    # Near start, middle, near end

# CIG locations for QC (inline, crossline pairs)
CIG_LOCATIONS = [
    (150, 150),  # Near center-left
    (255, 213),  # Center
    (350, 280),  # Center-right
    (200, 100),  # Upper region
    (300, 320),  # Lower region
]


# =============================================================================
# Helper Functions
# =============================================================================

def load_migration_volume(bin_path: Path) -> tuple[np.ndarray, dict]:
    """Load migrated volume from zarr."""
    zarr_path = bin_path / "migrated_stack.zarr"
    store = zarr.storage.LocalStore(str(zarr_path))
    z = zarr.open_array(store=store, mode='r')

    data = np.asarray(z)
    attrs = dict(z.attrs)

    return data, attrs


def get_bin_stats(bin_num: int, input_dir: Path) -> dict:
    """Get statistics for an offset bin."""
    headers_path = input_dir / f"offset_bin_{bin_num:02d}" / "headers.parquet"

    stats = {
        'bin_num': bin_num,
        'n_traces': 0,
        'mean_offset': 0.0,
        'min_offset': 0.0,
        'max_offset': 0.0,
    }

    if headers_path.exists():
        df = pl.read_parquet(headers_path)
        stats['n_traces'] = len(df)
        if 'offset' in df.columns and len(df) > 0:
            stats['mean_offset'] = float(df['offset'].mean())
            stats['min_offset'] = float(df['offset'].min())
            stats['max_offset'] = float(df['offset'].max())

    return stats


def plot_seismic_slice(data: np.ndarray, x_axis: np.ndarray, y_axis: np.ndarray,
                       title: str, xlabel: str, ylabel: str,
                       clip_percentile: float = 99, cmap: str = 'gray') -> BytesIO:
    """Create a seismic slice plot and return as bytes."""
    fig, ax = plt.subplots(figsize=(10, 8))

    vmax = np.percentile(np.abs(data), clip_percentile)
    if vmax == 0:
        vmax = 1.0
    vmin = -vmax

    extent = [x_axis[0], x_axis[-1], y_axis[-1], y_axis[0]]
    ax.imshow(data.T, aspect='auto', cmap=cmap, vmin=vmin, vmax=vmax,
              extent=extent, interpolation='bilinear')
    ax.set_xlabel(xlabel, fontsize=12)
    ax.set_ylabel(ylabel, fontsize=12)
    ax.set_title(title, fontsize=14, fontweight='bold')

    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close(fig)

    return buf


def plot_time_slice(data: np.ndarray, il_axis: np.ndarray, xl_axis: np.ndarray,
                    time_ms: float, clip_percentile: float = 99) -> BytesIO:
    """Create a time slice plot and return as bytes."""
    fig, ax = plt.subplots(figsize=(10, 8))

    vmax = np.percentile(np.abs(data), clip_percentile)
    if vmax == 0:
        vmax = 1.0
    vmin = -vmax

    extent = [il_axis[0], il_axis[-1], xl_axis[-1], xl_axis[0]]
    ax.imshow(data.T, aspect='auto', cmap='seismic', vmin=vmin, vmax=vmax,
              extent=extent, interpolation='bilinear')
    ax.set_xlabel('Inline', fontsize=12)
    ax.set_ylabel('Crossline', fontsize=12)
    ax.set_title(f'Time Slice @ {time_ms:.0f} ms', fontsize=14, fontweight='bold')

    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close(fig)

    return buf


def create_title_slide(prs: Presentation, title: str, subtitle: str):
    """Add a title slide to the presentation."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER


def create_stats_slide(prs: Presentation, stats: dict, migration_size: str):
    """Add a statistics slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"Offset Bin {stats['bin_num']:02d} - Statistics"
    p.font.size = Pt(32)
    p.font.bold = True

    # Stats content
    content = f"""
Input Data:
  • Number of traces: {stats['n_traces']:,}
  • Offset range: {stats['min_offset']:.0f} - {stats['max_offset']:.0f} m
  • Mean offset: {stats['mean_offset']:.0f} m

Migration Output:
  • Output size: {migration_size}
  • Grid: 511 IL x 427 XL x 1001 samples
  • Sample rate: 2 ms
  • Time range: 0 - 2000 ms
"""

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = content
    p2.font.size = Pt(20)


def create_image_slide(prs: Presentation, image_buf: BytesIO, title: str):
    """Add a slide with a single image."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Image
    slide.shapes.add_picture(image_buf, Inches(0.5), Inches(0.9), width=Inches(9))


def create_three_image_slide(prs: Presentation, images: list[BytesIO],
                              titles: list[str], main_title: str):
    """Add a slide with three images side by side."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.4), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = main_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Three images
    img_width = Inches(3.1)
    img_height = Inches(5.5)

    for i, (img_buf, title) in enumerate(zip(images, titles)):
        left = Inches(0.2 + i * 3.3)
        slide.shapes.add_picture(img_buf, left, Inches(0.7), width=img_width)

        # Subtitle under each image
        txBox2 = slide.shapes.add_textbox(left, Inches(6.3), img_width, Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER


# =============================================================================
# Main QC Functions
# =============================================================================

def create_bin_pptx(bin_num: int, migration_dir: Path, input_dir: Path,
                    output_dir: Path) -> Path:
    """Create a PPTX presentation for a single offset bin."""
    bin_path = migration_dir / f"migration_bin_{bin_num:02d}"
    zarr_path = bin_path / "migrated_stack.zarr"

    if not zarr_path.exists():
        print(f"  Skip bin {bin_num:02d}: no migrated data")
        return None

    print(f"  Processing bin {bin_num:02d}...")

    # Load data
    data, attrs = load_migration_volume(bin_path)
    stats = get_bin_stats(bin_num, input_dir)

    # Get migration size
    import subprocess
    size_result = subprocess.run(['du', '-sh', str(zarr_path)],
                                  capture_output=True, text=True)
    migration_size = size_result.stdout.split()[0] if size_result.returncode == 0 else "N/A"

    # Create axes
    nx, ny, nt = data.shape
    il_axis = np.arange(1, nx + 1)
    xl_axis = np.arange(1, ny + 1)
    dt_ms = attrs.get('dt_ms', 2.0)
    t_axis = np.arange(nt) * dt_ms

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        f"PSTM Migration QC",
        f"Offset Bin {bin_num:02d} - Mean Offset: {stats['mean_offset']:.0f} m"
    )

    # Stats slide
    create_stats_slide(prs, stats, migration_size)

    # Inline slices
    inline_images = []
    inline_titles = []
    for il_idx in INLINE_POSITIONS:
        if il_idx < nx:
            slice_data = data[il_idx, :, :]
            img = plot_seismic_slice(slice_data, xl_axis, t_axis,
                                     f'Inline {il_idx + 1}', 'Crossline', 'Time (ms)')
            inline_images.append(img)
            inline_titles.append(f'IL {il_idx + 1}')

    if inline_images:
        create_three_image_slide(prs, inline_images, inline_titles, "Inline Sections")

    # Crossline slices
    xline_images = []
    xline_titles = []
    for xl_idx in XLINE_POSITIONS:
        if xl_idx < ny:
            slice_data = data[:, xl_idx, :]
            img = plot_seismic_slice(slice_data, il_axis, t_axis,
                                     f'Crossline {xl_idx + 1}', 'Inline', 'Time (ms)')
            xline_images.append(img)
            xline_titles.append(f'XL {xl_idx + 1}')

    if xline_images:
        create_three_image_slide(prs, xline_images, xline_titles, "Crossline Sections")

    # Time slices
    time_images = []
    time_titles = []
    for time_ms in TIME_SLICES_MS:
        t_idx = int(time_ms / dt_ms)
        if t_idx < nt:
            slice_data = data[:, :, t_idx]
            img = plot_time_slice(slice_data, il_axis, xl_axis, time_ms)
            time_images.append(img)
            time_titles.append(f'{time_ms} ms')

    if time_images:
        create_three_image_slide(prs, time_images, time_titles, "Time Slices")

    # Save presentation
    output_path = output_dir / f"PSTM_QC_Bin_{bin_num:02d}.pptx"
    prs.save(str(output_path))

    return output_path


def create_summary_pptx(migration_dir: Path, input_dir: Path,
                        output_dir: Path, bins: list[int]) -> Path:
    """Create a summary PPTX with all bins comparison."""
    print("Creating summary presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        "PSTM Migration Summary",
        f"{len(bins)} Offset Bins - {datetime.now().strftime('%Y-%m-%d')}"
    )

    # Bin statistics table slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Offset Bin Statistics"
    p.font.size = Pt(28)
    p.font.bold = True

    # Stats text
    stats_text = "Bin | Traces     | Offset Range (m) | Migration Size\n"
    stats_text += "-" * 55 + "\n"

    for bin_num in bins[:20]:  # First 20 bins
        stats = get_bin_stats(bin_num, input_dir)
        zarr_path = migration_dir / f"migration_bin_{bin_num:02d}" / "migrated_stack.zarr"
        if zarr_path.exists():
            import subprocess
            size_result = subprocess.run(['du', '-sh', str(zarr_path)],
                                          capture_output=True, text=True)
            size = size_result.stdout.split()[0] if size_result.returncode == 0 else "N/A"
        else:
            size = "N/A"

        stats_text += f" {bin_num:02d} | {stats['n_traces']:>10,} | {stats['min_offset']:>5.0f} - {stats['max_offset']:<5.0f} | {size}\n"

    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = stats_text
    p2.font.size = Pt(12)
    p2.font.name = 'Courier New'

    # Continue with bins 20-39 if present
    if len(bins) > 20:
        slide = prs.slides.add_slide(slide_layout)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Offset Bin Statistics (continued)"
        p.font.size = Pt(28)
        p.font.bold = True

        stats_text = "Bin | Traces     | Offset Range (m) | Migration Size\n"
        stats_text += "-" * 55 + "\n"

        for bin_num in bins[20:]:
            stats = get_bin_stats(bin_num, input_dir)
            zarr_path = migration_dir / f"migration_bin_{bin_num:02d}" / "migrated_stack.zarr"
            if zarr_path.exists():
                import subprocess
                size_result = subprocess.run(['du', '-sh', str(zarr_path)],
                                              capture_output=True, text=True)
                size = size_result.stdout.split()[0] if size_result.returncode == 0 else "N/A"
            else:
                size = "N/A"

            stats_text += f" {bin_num:02d} | {stats['n_traces']:>10,} | {stats['min_offset']:>5.0f} - {stats['max_offset']:<5.0f} | {size}\n"

        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = stats_text
        p2.font.size = Pt(12)
        p2.font.name = 'Courier New'

    # Save
    output_path = output_dir / "PSTM_QC_Summary.pptx"
    prs.save(str(output_path))

    return output_path


# =============================================================================
# CIG Functions
# =============================================================================

def extract_cig_at_location(migration_dir: Path, il_idx: int, xl_idx: int,
                            bins: list[int]) -> tuple[np.ndarray, np.ndarray]:
    """Extract CIG traces at a specific location from all bins."""
    traces = []
    offsets = []

    for bin_num in bins:
        bin_path = migration_dir / f"migration_bin_{bin_num:02d}"
        zarr_path = bin_path / "migrated_stack.zarr"

        if not zarr_path.exists():
            continue

        store = zarr.storage.LocalStore(str(zarr_path))
        z = zarr.open_array(store=store, mode='r')

        # Get trace at location
        if il_idx < z.shape[0] and xl_idx < z.shape[1]:
            trace = z[il_idx, xl_idx, :]
            traces.append(trace)

            # Estimate offset from bin number (50m bins)
            offset = bin_num * 50 + 25
            offsets.append(offset)

    return np.array(traces), np.array(offsets)


def plot_cig(traces: np.ndarray, offsets: np.ndarray, il: int, xl: int,
             dt_ms: float = 2.0, clip_percentile: float = 99) -> BytesIO:
    """Create a CIG plot."""
    fig, ax = plt.subplots(figsize=(8, 10))

    n_traces, nt = traces.shape
    t_axis = np.arange(nt) * dt_ms

    vmax = np.percentile(np.abs(traces), clip_percentile)
    if vmax == 0:
        vmax = 1.0

    extent = [offsets[0], offsets[-1], t_axis[-1], t_axis[0]]
    ax.imshow(traces.T, aspect='auto', cmap='gray', vmin=-vmax, vmax=vmax,
              extent=extent, interpolation='bilinear')

    ax.set_xlabel('Offset (m)', fontsize=12)
    ax.set_ylabel('Time (ms)', fontsize=12)
    ax.set_title(f'CIG at IL={il}, XL={xl}', fontsize=14, fontweight='bold')

    plt.tight_layout()

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    buf.seek(0)
    plt.close(fig)

    return buf


def create_cig_pptx(migration_dir: Path, output_dir: Path,
                    bins: list[int], locations: list[tuple]) -> Path:
    """Create a PPTX presentation for CIG QC."""
    print("Creating CIG presentation...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    create_title_slide(
        prs,
        "Common Image Gathers (CIG) QC",
        f"{len(locations)} Locations - {datetime.now().strftime('%Y-%m-%d')}"
    )

    # CIG slides
    for il, xl in locations:
        print(f"    Extracting CIG at IL={il}, XL={xl}...")

        traces, offsets = extract_cig_at_location(migration_dir, il-1, xl-1, bins)

        if len(traces) > 0:
            img = plot_cig(traces, offsets, il, xl)
            create_image_slide(prs, img, f"CIG at Inline {il}, Crossline {xl}")

    # Save
    output_path = output_dir / "PSTM_CIG_QC.pptx"
    prs.save(str(output_path))

    return output_path


# =============================================================================
# Main Entry Point
# =============================================================================

def main():
    parser = argparse.ArgumentParser(description="Generate PSTM QC PPTX reports")
    parser.add_argument("--bins", type=str, default="0-37",
                        help="Bins to process (e.g., '0-37' or 'all')")
    parser.add_argument("--migration-dir", type=Path, default=MIGRATION_DIR,
                        help="Migration output directory")
    parser.add_argument("--input-dir", type=Path, default=INPUT_DIR,
                        help="Input data directory")
    parser.add_argument("--output-dir", type=Path, default=OUTPUT_DIR,
                        help="Output directory for reports")
    parser.add_argument("--skip-individual", action="store_true",
                        help="Skip individual bin presentations")
    parser.add_argument("--cig-only", action="store_true",
                        help="Only generate CIG presentation")

    args = parser.parse_args()

    # Parse bins
    if args.bins.lower() == "all":
        bins = list(range(40))
    elif "-" in args.bins:
        start, end = map(int, args.bins.split("-"))
        bins = list(range(start, end + 1))
    else:
        bins = [int(b) for b in args.bins.split(",")]

    # Filter to existing bins
    bins = [b for b in bins if (args.migration_dir / f"migration_bin_{b:02d}" / "migrated_stack.zarr").exists()]

    print("=" * 70)
    print("PSTM Migration QC Report Generator")
    print("=" * 70)
    print(f"Migration directory: {args.migration_dir}")
    print(f"Output directory: {args.output_dir}")
    print(f"Bins to process: {len(bins)} ({min(bins)}-{max(bins)})")
    print()

    # Create output directory
    args.output_dir.mkdir(parents=True, exist_ok=True)

    if not args.cig_only:
        # Individual bin presentations
        if not args.skip_individual:
            print("[1] Creating individual bin presentations...")
            for bin_num in bins:
                result = create_bin_pptx(bin_num, args.migration_dir, args.input_dir,
                                         args.output_dir)
                if result:
                    print(f"      -> {result.name}")
            print()

        # Summary presentation
        print("[2] Creating summary presentation...")
        summary_path = create_summary_pptx(args.migration_dir, args.input_dir,
                                           args.output_dir, bins)
        print(f"    -> {summary_path.name}")
        print()

    # CIG presentation
    print("[3] Creating CIG presentation...")
    cig_path = create_cig_pptx(args.migration_dir, args.output_dir, bins, CIG_LOCATIONS)
    print(f"    -> {cig_path.name}")
    print()

    print("=" * 70)
    print("QC Report Generation Complete")
    print("=" * 70)
    print(f"Output directory: {args.output_dir}")
    print()

    return 0


if __name__ == "__main__":
    sys.exit(main())
